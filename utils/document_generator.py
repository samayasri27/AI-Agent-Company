"""
Document Generator Module
Generates professional documents in various formats (DOCX, PPTX, XLSX, PDF)
"""
import io
import os
from datetime import datetime
from typing import Dict, List, Any, Optional

try:
    from docx import Document
    from docx.shared import Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment
    from openpyxl.chart import BarChart, Reference
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False


class DocumentGenerator:
    """Professional document generator for various formats"""
    
    def __init__(self):
        self.available_formats = []
        if DOCX_AVAILABLE:
            self.available_formats.append('docx')
        if PPTX_AVAILABLE:
            self.available_formats.append('pptx')
        if XLSX_AVAILABLE:
            self.available_formats.append('xlsx')
        if PDF_AVAILABLE:
            self.available_formats.append('pdf')
    
    def create_docx(self, content: Dict[str, Any], template: str = None) -> bytes:
        """Create DOCX document"""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx not available. Install with: pip install python-docx")
        
        doc = Document()
        
        # Add title
        if 'title' in content:
            title = doc.add_heading(content['title'], 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add metadata
        if 'author' in content:
            doc.core_properties.author = content['author']
        if 'subject' in content:
            doc.core_properties.subject = content['subject']
        
        doc.core_properties.created = datetime.now()
        
        # Add executive summary
        if 'executive_summary' in content:
            doc.add_heading('Executive Summary', level=1)
            doc.add_paragraph(content['executive_summary'])
            doc.add_page_break()
        
        # Add sections
        if 'sections' in content:
            for section in content['sections']:
                doc.add_heading(section.get('title', 'Section'), level=1)
                
                if 'content' in section:
                    doc.add_paragraph(section['content'])
                
                if 'subsections' in section:
                    for subsection in section['subsections']:
                        doc.add_heading(subsection.get('title', 'Subsection'), level=2)
                        if 'content' in subsection:
                            doc.add_paragraph(subsection['content'])
                
                if 'bullet_points' in section:
                    for point in section['bullet_points']:
                        doc.add_paragraph(point, style='List Bullet')
                
                if 'table' in section:
                    self._add_table_to_docx(doc, section['table'])
        
        # Add findings/recommendations
        if 'findings' in content:
            doc.add_heading('Key Findings', level=1)
            for finding in content['findings']:
                doc.add_paragraph(finding, style='List Bullet')
        
        if 'recommendations' in content:
            doc.add_heading('Recommendations', level=1)
            for rec in content['recommendations']:
                doc.add_paragraph(rec, style='List Bullet')
        
        # Save to bytes
        doc_io = io.BytesIO()
        doc.save(doc_io)
        doc_io.seek(0)
        return doc_io.getvalue()
    
    def create_pptx(self, slides: List[Dict[str, Any]], template: str = None) -> bytes:
        """Create PPTX presentation"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx not available. Install with: pip install python-pptx")
        
        prs = Presentation()
        
        for slide_data in slides:
            # Choose slide layout
            layout_type = slide_data.get('layout', 'title_and_content')
            if layout_type == 'title_slide':
                slide_layout = prs.slide_layouts[0]  # Title slide
            elif layout_type == 'title_only':
                slide_layout = prs.slide_layouts[5]  # Title only
            else:
                slide_layout = prs.slide_layouts[1]  # Title and content
            
            slide = prs.slides.add_slide(slide_layout)
            
            # Add title
            if 'title' in slide_data:
                title = slide.shapes.title
                title.text = slide_data['title']
            
            # Add content
            if 'content' in slide_data and hasattr(slide.shapes, 'placeholders'):
                if len(slide.shapes.placeholders) > 1:
                    content_placeholder = slide.shapes.placeholders[1]
                    if hasattr(content_placeholder, 'text_frame'):
                        tf = content_placeholder.text_frame
                        tf.text = slide_data['content']
            
            # Add bullet points
            if 'bullet_points' in slide_data and hasattr(slide.shapes, 'placeholders'):
                if len(slide.shapes.placeholders) > 1:
                    content_placeholder = slide.shapes.placeholders[1]
                    if hasattr(content_placeholder, 'text_frame'):
                        tf = content_placeholder.text_frame
                        tf.clear()
                        for i, point in enumerate(slide_data['bullet_points']):
                            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                            p.text = point
                            p.level = 0
            
            # Add image if specified
            if 'image' in slide_data:
                # Add image placeholder functionality
                pass
        
        # Save to bytes
        prs_io = io.BytesIO()
        prs.save(prs_io)
        prs_io.seek(0)
        return prs_io.getvalue()
    
    def create_xlsx(self, data: Dict[str, Any], sheets: List[str] = None) -> bytes:
        """Create XLSX spreadsheet"""
        if not XLSX_AVAILABLE:
            raise ImportError("openpyxl not available. Install with: pip install openpyxl")
        
        wb = Workbook()
        
        # Remove default sheet if we have custom sheets
        if sheets and 'Sheet' in [ws.title for ws in wb.worksheets]:
            wb.remove(wb['Sheet'])
        
        sheets_to_create = sheets or ['Data']
        
        for sheet_name in sheets_to_create:
            if sheet_name not in [ws.title for ws in wb.worksheets]:
                ws = wb.create_sheet(title=sheet_name)
            else:
                ws = wb[sheet_name]
            
            # Add data for this sheet
            sheet_data = data.get(sheet_name.lower(), data.get('data', {}))
            
            if 'headers' in sheet_data and 'rows' in sheet_data:
                # Add headers
                headers = sheet_data['headers']
                for col, header in enumerate(headers, 1):
                    cell = ws.cell(row=1, column=col, value=header)
                    cell.font = Font(bold=True)
                    cell.fill = PatternFill(start_color="CCCCCC", end_color="CCCCCC", fill_type="solid")
                
                # Add data rows
                for row_idx, row_data in enumerate(sheet_data['rows'], 2):
                    for col_idx, value in enumerate(row_data, 1):
                        ws.cell(row=row_idx, column=col_idx, value=value)
                
                # Auto-adjust column widths
                for column in ws.columns:
                    max_length = 0
                    column_letter = column[0].column_letter
                    for cell in column:
                        try:
                            if len(str(cell.value)) > max_length:
                                max_length = len(str(cell.value))
                        except:
                            pass
                    adjusted_width = min(max_length + 2, 50)
                    ws.column_dimensions[column_letter].width = adjusted_width
            
            # Add chart if specified
            if 'chart' in sheet_data:
                self._add_chart_to_xlsx(ws, sheet_data['chart'])
        
        # Save to bytes
        wb_io = io.BytesIO()
        wb.save(wb_io)
        wb_io.seek(0)
        return wb_io.getvalue()
    
    def create_pdf(self, content: Dict[str, Any], template: str = None) -> bytes:
        """Create PDF document"""
        if not PDF_AVAILABLE:
            raise ImportError("reportlab not available. Install with: pip install reportlab")
        
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            spaceAfter=30,
            alignment=1  # Center
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=16,
            spaceAfter=12,
            textColor=colors.darkblue
        )
        
        # Add title
        if 'title' in content:
            story.append(Paragraph(content['title'], title_style))
            story.append(Spacer(1, 12))
        
        # Add metadata
        if 'author' in content or 'date' in content:
            meta_text = []
            if 'author' in content:
                meta_text.append(f"Author: {content['author']}")
            if 'date' in content:
                meta_text.append(f"Date: {content['date']}")
            else:
                meta_text.append(f"Date: {datetime.now().strftime('%Y-%m-%d')}")
            
            story.append(Paragraph(" | ".join(meta_text), styles['Normal']))
            story.append(Spacer(1, 20))
        
        # Add executive summary
        if 'executive_summary' in content:
            story.append(Paragraph("Executive Summary", heading_style))
            story.append(Paragraph(content['executive_summary'], styles['Normal']))
            story.append(Spacer(1, 20))
        
        # Add sections
        if 'sections' in content:
            for section in content['sections']:
                story.append(Paragraph(section.get('title', 'Section'), heading_style))
                
                if 'content' in section:
                    story.append(Paragraph(section['content'], styles['Normal']))
                    story.append(Spacer(1, 12))
                
                if 'bullet_points' in section:
                    for point in section['bullet_points']:
                        story.append(Paragraph(f"• {point}", styles['Normal']))
                    story.append(Spacer(1, 12))
                
                if 'table' in section:
                    table_data = section['table']
                    if 'headers' in table_data and 'rows' in table_data:
                        data = [table_data['headers']] + table_data['rows']
                        table = Table(data)
                        table.setStyle(TableStyle([
                            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                            ('FONTSIZE', (0, 0), (-1, 0), 14),
                            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                            ('GRID', (0, 0), (-1, -1), 1, colors.black)
                        ]))
                        story.append(table)
                        story.append(Spacer(1, 12))
        
        # Add findings/recommendations
        if 'findings' in content:
            story.append(Paragraph("Key Findings", heading_style))
            for finding in content['findings']:
                story.append(Paragraph(f"• {finding}", styles['Normal']))
            story.append(Spacer(1, 12))
        
        if 'recommendations' in content:
            story.append(Paragraph("Recommendations", heading_style))
            for rec in content['recommendations']:
                story.append(Paragraph(f"• {rec}", styles['Normal']))
        
        # Build PDF
        doc.build(story)
        buffer.seek(0)
        return buffer.getvalue()
    
    def _add_table_to_docx(self, doc, table_data: Dict[str, Any]):
        """Add table to DOCX document"""
        if 'headers' in table_data and 'rows' in table_data:
            headers = table_data['headers']
            rows = table_data['rows']
            
            table = doc.add_table(rows=1, cols=len(headers))
            table.style = 'Table Grid'
            
            # Add headers
            hdr_cells = table.rows[0].cells
            for i, header in enumerate(headers):
                hdr_cells[i].text = str(header)
            
            # Add data rows
            for row_data in rows:
                row_cells = table.add_row().cells
                for i, value in enumerate(row_data):
                    row_cells[i].text = str(value)
    
    def _add_chart_to_xlsx(self, ws, chart_data: Dict[str, Any]):
        """Add chart to XLSX worksheet"""
        chart_type = chart_data.get('type', 'bar')
        
        if chart_type == 'bar':
            chart = BarChart()
            chart.title = chart_data.get('title', 'Chart')
            chart.y_axis.title = chart_data.get('y_axis', 'Values')
            chart.x_axis.title = chart_data.get('x_axis', 'Categories')
            
            # Assume data is in the worksheet
            data_range = chart_data.get('data_range', 'A1:B10')
            data = Reference(ws, range_string=data_range)
            chart.add_data(data, titles_from_data=True)
            
            # Add chart to worksheet
            ws.add_chart(chart, chart_data.get('position', 'D2'))
    
    def get_available_formats(self) -> List[str]:
        """Get list of available document formats"""
        return self.available_formats
    
    def check_dependencies(self) -> Dict[str, bool]:
        """Check which document generation dependencies are available"""
        return {
            'docx': DOCX_AVAILABLE,
            'pptx': PPTX_AVAILABLE,
            'xlsx': XLSX_AVAILABLE,
            'pdf': PDF_AVAILABLE
        }